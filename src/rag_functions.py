from io import BytesIO
from PyPDF2 import PdfReader
from pptx import Presentation
import logging
import os
from dotenv import load_dotenv
from werkzeug.utils import secure_filename
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from langchain.text_splitter import RecursiveCharacterTextSplitter

import os

# Handling vector database
from langchain import PromptTemplate

# Add OpenAI library
import openai

# Setting up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure OpenAI API using Azure OpenAI
openai.api_key = os.getenv("API_KEY")
openai.api_base = os.getenv("ENDPOINT")
openai.api_type = "azure"  # Necessary for using the OpenAI library with Azure OpenAI
openai.api_version = os.getenv("OPENAI_API_VERSION")  # Latest / target version of the API

from langchain.embeddings import OpenAIEmbeddings

# OpenAI Settings
model_deployment = "text-embedding-ada-002"
# SDK calls this "engine", but naming it "deployment_name" for clarity

model_name = "text-embedding-ada-002"

# Embeddings
openai_embeddings: OpenAIEmbeddings = OpenAIEmbeddings(
    openai_api_version = os.getenv("OPENAI_API_VERSION"), openai_api_key = os.getenv("API_KEY"),
    openai_api_base = os.getenv("ENDPOINT"), openai_api_type = "azure"
)


# Load environment variables from .env file
load_dotenv()

# Allowed file types
allowed_files_list = ["pdf", "txt", "pptx"]

def allowed_files(filename):
    '''
    Returns True if the file type is in the allowed file list
    '''
    return "." in filename and filename.rsplit(".", 1)[1].lower() in allowed_files_list

def file_check_num(uploaded_file):
    '''
    Returns the number of pages (for PDFs), slides (for PPTX), or lines (for TXT) in the file
    '''
    file_ext = uploaded_file.name.rsplit(".", 1)[1].lower()  # Extract the file extension
    try:
        if file_ext == "pdf":
            pdf_bytes = BytesIO(uploaded_file.read())
            pdf_reader = PdfReader(pdf_bytes)
            uploaded_file.seek(0)  # Reset file pointer after reading
            return len(pdf_reader.pages)
        
        elif file_ext == "pptx":
            pptx_bytes = BytesIO(uploaded_file.read())
            pptx = Presentation(pptx_bytes)
            uploaded_file.seek(0)
            return len(pptx.slides)
        
        elif file_ext == "txt":
            num = len(uploaded_file.read().decode("utf-8").splitlines())
            uploaded_file.seek(0)
            return num
        else:
            logger.error(f"Unsupported file extension: {file_ext}")
            return -1
    except Exception as e:
        logger.error(f"Error checking file '{uploaded_file.name}': {e}")
        return -1



def chunk_document(text, chunk_size=1000, chunk_overlap=300):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n", " ", "?", ".", "!"]
    )
    chunks = text_splitter.split_text(text)
    return chunks

def extract_contents_from_doc(files, temp_dir):
    """
    Azure Document Intelligence
    Args: 
        files (uploaded by the user): List of uploaded files to process.
        temp_dir (str): Directory path to store the extracted contents.
    
    Returns: 
        List of file paths where the extracted content is stored.
    """
    # Constants for Azure Document Intelligence
    DI_ENDPOINT = os.getenv("DOCUMENT_INTELLIGENCE_ENDPOINT")
    DOCUMENT_INTELLIGENCE_KEY = os.getenv('DOCUMENT_INTELLIGENCE_SUBSCRIPTION_KEY')

    if not DI_ENDPOINT or not DOCUMENT_INTELLIGENCE_KEY:
        logger.error("Azure Document Intelligence credentials are missing.")
        return []

    document_intelligence_client = DocumentAnalysisClient(
        endpoint=DI_ENDPOINT,
        credential=AzureKeyCredential(DOCUMENT_INTELLIGENCE_KEY)
    )

    # Ensure the temporary directory exists
    os.makedirs(temp_dir, exist_ok=True)
    logger.info(f"Temporary directory '{temp_dir}' is ready.")

    extracted_file_paths = []

    for file in files:
        filename = secure_filename(file.name)
        base, ext = os.path.splitext(filename)
        ext = ext.lower()  # Normalize to lowercase for easier checking
        
        try:
            if ext == '.pdf':
                # Extract content using Azure Document Intelligence for PDF
                file_content = file.read()
                logger.info(f"Processing PDF file: {file.name}")
                extract = document_intelligence_client.begin_analyze_document("prebuilt-read", file_content)
                result = extract.result()
                
                # Extract text from each page
                extracted_content = ""
                for page in result.pages:
                    for line in page.lines:
                        extracted_content += line.content + "\n"
                
            elif ext == '.txt':
                # Directly read .txt files
                logger.info(f"Processing TXT file: {file.name}")
                extracted_content = file.read().decode('utf-8')
                
            elif ext == '.pptx':
                # Extract content from .pptx using python-pptx
                logger.info(f"Processing PPTX file: {file.name}")
                extracted_content = ""
                presentation = Presentation(file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_content += shape.text + "\n"
            
            else:
                logger.warning(f"Unsupported file type: {file.name}")
                continue  # Skip unsupported file types
            
            # Define path to save extracted content
            extracted_filename = f"{base}_extracted.txt"  # Save as .txt for easier reading
            file_path = os.path.join(temp_dir, extracted_filename)
            
            # Save the extracted content to a file
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(extracted_content)
            
            logger.info(f"Extracted content saved to: {file_path}")
            extracted_file_paths.append(file_path)

        except Exception as e:
            logger.error(f"Error processing file '{file.name}': {e}")
            continue  # Proceed with the next file in case of an error

    return extracted_file_paths

def conversation_history_prompt(history, question):
    # Define the template string for summarizing conversation history
    template_summary = """
    "Given a chat history (delimited by <hs></hs>) and the latest user question \
    which might reference context in the chat history, formulate a standalone question \
    which can be understood without the chat history. Do NOT answer the question, \
    just reformulate it if needed and otherwise return it as is.
    ------
    <hs>
    {history}
    </hs>
    ------
    Question: {question}
    Summary:
    """

    # Create a PromptTemplate object
    prompt = PromptTemplate(
        input_variables=["history", "question"],
        template=template_summary,
    )

    return prompt.format(history=history, question=question)

def get_conversation_summary(history, question):
    # Get the conversation summary prompt
    formatted_prompt = conversation_history_prompt(history, question)

    # Query the Azure OpenAI LLM with the formatted prompt
    response = openai.ChatCompletion.create(
        engine="Voicetask",  # Replace with your Azure OpenAI deployment name
        # prompt=formatted_prompt,
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": formatted_prompt}
        ],
        # max_tokens=50,
        temperature=0.5
    )
    
    # Extract and return the summary from the response
    return response.choices[0].message['content']
